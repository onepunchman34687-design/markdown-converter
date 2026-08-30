"""
image_extract.py
Pulls embedded images out of PPTX, DOCX, and PDF files as raw bytes,
then attaches a text description to each one (skipping tiny icons/bullets
cheaply via Pillow before paying for any vision call). The description
is what makes these embedded images actually readable by a text-only LLM
once the markdown is pushed — same fix as the standalone image path.
"""
import io
from PIL import Image
from image_describe import describe_image, MIN_DESCRIBE_DIMENSION


def _is_too_small(image_bytes):
    """Cheap local check so we never pay for a vision call on icons/bullets."""
    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            return max(img.size) < MIN_DESCRIBE_DIMENSION
    except Exception:
        return False  # if we can't even open it, let describe_image() decide


def extract_images(file_path, suffix):
    """
    Dispatches to the right extractor based on file extension.
    Returns a list of dicts:
        [{'filename': ..., 'data': bytes, 'content_type': ..., 'description': str|None}, ...]
    Never raises — if extraction fails or the format isn't supported, returns [].
    """
    suffix = (suffix or '').lower()
    try:
        if suffix == '.pptx':
            images = _extract_pptx(file_path)
        elif suffix == '.docx':
            images = _extract_docx(file_path)
        elif suffix == '.pdf':
            images = _extract_pdf(file_path)
        else:
            images = []
    except Exception as e:
        print(f"Image extraction skipped ({suffix}): {e}")
        return []

    for img in images:
        if _is_too_small(img['data']):
            img['description'] = None
            continue
        img['description'] = describe_image(img['data'], img['content_type'])

    return images


def _extract_pptx(file_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    images = []
    prs = Presentation(file_path)
    idx = 1
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext or 'png'
                images.append({
                    'filename': f'image_{idx}.{ext}',
                    'data': image.blob,
                    'content_type': image.content_type or f'image/{ext}'
                })
                idx += 1
    return images

def _extract_docx(file_path):
    from docx import Document
    images = []
    doc = Document(file_path)
    idx = 1
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_part = rel.target_part
            content_type = image_part.content_type or 'image/png'
            ext = content_type.split('/')[-1]
            images.append({
                'filename': f'image_{idx}.{ext}',
                'data': image_part.blob,
                'content_type': content_type
            })
            idx += 1
    return images

def _extract_pdf(file_path):
    import fitz  # PyMuPDF
    images = []
    doc = fitz.open(file_path)
    idx = 1
    for page_index in range(len(doc)):
        page = doc[page_index]
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            ext = base_image.get('ext', 'png')
            images.append({
                'filename': f'image_{idx}.{ext}',
                'data': base_image['image'],
                'content_type': f'image/{ext}'
            })
            idx += 1
    doc.close()
    return images